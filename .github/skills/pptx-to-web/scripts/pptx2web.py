#!/usr/bin/env python3
"""pptx2web — convert a .pptx into a static HTML deck for GitHub Pages.

Pipeline: soffice (pptx->pdf) -> PyMuPDF (pdf page -> PNG) for exact look,
python-pptx for video overlay coords, zip extract for media. Emits docs/.

Usage: python3 pptx2web.py <input.pptx> [--out docs] [--scale 2.0]
"""
import argparse, json, os, shutil, subprocess, sys, tempfile, zipfile
from pathlib import Path

VIDEO_EXT = {".mp4", ".m4v", ".mov", ".webm", ".ogv", ".avi", ".wmv", ".mkv"}
AUDIO_EXT = {".mp3", ".m4a", ".wav", ".aac", ".ogg", ".wma"}
BROWSER_OK = {".mp4", ".m4v", ".webm", ".ogv", ".mp3", ".m4a", ".wav", ".aac", ".ogg"}


def find_soffice():
    for c in ("soffice", "libreoffice", "/Applications/LibreOffice.app/Contents/MacOS/soffice"):
        if shutil.which(c) or os.path.exists(c):
            return c
    sys.exit("LibreOffice (soffice) not found. Install it to render slides.")


def render_pngs(pptx, outdir, scale):
    import fitz  # PyMuPDF
    soffice = find_soffice()
    with tempfile.TemporaryDirectory() as tmp:
        subprocess.run([soffice, "--headless", "--convert-to", "pdf", "--outdir", tmp, pptx],
                       check=True, capture_output=True)
        pdfs = list(Path(tmp).glob("*.pdf"))
        if not pdfs:
            sys.exit(f"Conversion failed: LibreOffice produced no PDF for {pptx}")
        doc = fitz.open(pdfs[0])
        n = doc.page_count
        for i in range(n):
            pix = doc[i].get_pixmap(matrix=fitz.Matrix(scale, scale))
            pix.save(str(outdir / f"slide-{i+1:03d}.png"))
        return n


def online_videos(pptx):
    """slide_index -> [embed urls] from external relationships (YouTube/Vimeo/etc)."""
    import re
    out = {}
    with zipfile.ZipFile(pptx) as z:
        for nm in z.namelist():
            m = re.match(r"ppt/slides/_rels/slide(\d+)\.xml\.rels", nm)
            if not m:
                continue
            idx = int(m.group(1)) - 1
            data = z.read(nm).decode("utf8", "ignore")
            urls = [u.replace("&amp;", "&") for u in re.findall(r'Target="([^"]+)"[^>]*External', data) if u.startswith("http")]
            urls = [u for u in urls if any(k in u for k in ("youtube", "youtu.be", "vimeo", "embed", ".mp4", ".webm"))]
            if urls:
                out[idx] = urls
    return out


def media_overlays(pptx):
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    p = Presentation(pptx)
    W, H = p.slide_width, p.slide_height
    by_slide = {}
    for i, s in enumerate(p.slides):
        items = []
        for sh in s.shapes:
            if sh.shape_type == MSO_SHAPE_TYPE.MEDIA and None not in (sh.left, sh.top, sh.width, sh.height):
                items.append({"left": max(0, sh.left/W*100), "top": max(0, sh.top/H*100),
                              "width": min(100, sh.width/W*100), "height": min(100, sh.height/H*100),
                              "name": sh.name})
        if items:
            by_slide[i] = items
    return by_slide, W/H


def transcode(path):
    """Transcode a non-browser video to mp4 (h264/aac) if ffmpeg is present."""
    if not shutil.which("ffmpeg"):
        return path
    mp4 = path.with_suffix(".mp4")
    r = subprocess.run(["ffmpeg", "-y", "-i", str(path), "-c:v", "libx264", "-c:a", "aac",
                        "-movflags", "+faststart", str(mp4)], capture_output=True)
    return mp4 if r.returncode == 0 and mp4.exists() else path


def extract_media(pptx, mediadir):
    files = []
    with zipfile.ZipFile(pptx) as z:
        for nm in z.namelist():
            ext = os.path.splitext(nm)[1].lower()
            if nm.startswith("ppt/media/") and ext in VIDEO_EXT | AUDIO_EXT:
                dst = mediadir / os.path.basename(nm)
                with z.open(nm) as src, open(dst, "wb") as out:
                    shutil.copyfileobj(src, out)
                if ext in VIDEO_EXT and ext not in BROWSER_OK:
                    new = transcode(dst)
                    if new != dst:
                        dst.unlink(missing_ok=True)
                        dst = new
                files.append(dst.name)
    files.sort()
    return files


def assign(overlays, media):
    vids = [m for m in media if os.path.splitext(m)[1].lower() in VIDEO_EXT]
    vi = 0
    for sl in sorted(overlays):
        for ov in overlays[sl]:
            src = vids[vi % len(vids)] if vids else None
            ov["src"] = ("media/" + src) if src else None
            ov["ok"] = bool(src) and os.path.splitext(src)[1].lower() in BROWSER_OK
            vi += 1
    return overlays


def build(pptx, out, scale):
    out = Path(out)
    for d in (out, out/"slides", out/"media"):
        d.mkdir(parents=True, exist_ok=True)
    for f in (out/"slides").glob("*"):
        f.unlink()
    for f in (out/"media").glob("*"):
        f.unlink()
    n = render_pngs(pptx, out/"slides", scale)
    overlays, ar = media_overlays(pptx)
    media = extract_media(pptx, out/"media")
    overlays = assign(overlays, media)
    slides = [{"img": f"slides/slide-{i+1:03d}.png", "videos": overlays.get(i, [])} for i in range(n)]
    deck = {"title": Path(pptx).stem, "aspect": ar, "slides": slides}
    (out/"deck.json").write_text(json.dumps(deck, ensure_ascii=False, indent=2))
    (out/"index.html").write_text(HTML.replace("__DECK__", json.dumps(deck, ensure_ascii=False)))
    return n, sum(len(s["videos"]) for s in slides)


HTML = r"""<!doctype html><html lang="ko"><head><meta charset="utf-8"/>
<meta name="viewport" content="width=device-width,initial-scale=1"/>
<title>Deck</title><style>
*{box-sizing:border-box}html,body{margin:0;height:100%;background:#111;font-family:system-ui,sans-serif}
#stage{position:absolute;inset:0;display:flex;align-items:center;justify-content:center}
#slide{position:relative;background:#000}#slide img{display:block;width:100%;height:100%;object-fit:contain}
.ov{position:absolute}.ov video{width:100%;height:100%}
.ov .na{width:100%;height:100%;display:flex;align-items:center;justify-content:center;color:#fff;background:rgba(0,0,0,.6);text-align:center;font-size:13px}
.ov .na a{color:#58a6ff}
#bar{position:fixed;left:0;bottom:0;height:4px;background:#58a6ff;transition:width .2s}
#hud{position:fixed;right:12px;bottom:12px;color:#aaa;font-size:13px;background:rgba(0,0,0,.5);padding:4px 8px;border-radius:6px}
.nav{position:fixed;top:0;bottom:0;width:18%;cursor:pointer}#prev{left:0}#next{right:0}
</style></head><body>
<div id="stage"><div id="slide"></div></div>
<div id="prev" class="nav"></div><div id="next" class="nav"></div>
<div id="bar"></div><div id="hud"></div>
<script>
const D=__DECK__;let i=0;const sl=document.getElementById('slide');
function layout(){const ar=D.aspect||16/9,w=innerWidth,h=innerHeight;let cw=w,ch=w/ar;if(ch>h){ch=h;cw=h*ar;}sl.style.width=cw+'px';sl.style.height=ch+'px';}
function show(n){i=Math.max(0,Math.min(D.slides.length-1,n));const s=D.slides[i];
sl.innerHTML='<img src="'+s.img+'"/>'+s.videos.map(v=>{const st='left:'+v.left+'%;top:'+v.top+'%;width:'+v.width+'%;height:'+v.height+'%;';
return '<div class="ov" style="'+st+'">'+(v.ok?'<video src="'+v.src+'" controls playsinline></video>':v.src?'<div class="na">브라우저 미지원 코덱<br><a href="'+v.src+'">영상 다운로드</a></div>':'')+'</div>';}).join('');
document.getElementById('bar').style.width=((i+1)/D.slides.length*100)+'%';document.getElementById('hud').textContent=(i+1)+' / '+D.slides.length;layout();}
addEventListener('keydown',e=>{if(['ArrowRight','PageDown',' '].includes(e.key))show(i+1);if(['ArrowLeft','PageUp'].includes(e.key))show(i-1);if(e.key=='f')document.documentElement.requestFullscreen();});
document.getElementById('next').onclick=()=>show(i+1);document.getElementById('prev').onclick=()=>show(i-1);addEventListener('resize',layout);show(0);document.title=D.title;
</script></body></html>"""


if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("pptx"); ap.add_argument("--out", default="docs"); ap.add_argument("--scale", type=float, default=2.0)
    a = ap.parse_args()
    if not os.path.exists(a.pptx):
        sys.exit(f"Input not found: {a.pptx}")
    s, v = build(a.pptx, a.out, a.scale)
    print(f"OK: {s} slides, {v} video overlays -> {a.out}/")
