#!/usr/bin/env python3
"""pptx2web_native — rebuild a .pptx as a real HTML deck (content, not screenshots).

Extracts text runs (font, size, color, bold/italic, align), pictures, tables and
embedded/online videos from each slide and lays them out as absolutely-positioned
HTML/CSS so text stays selectable and the layout adapts to the web. Emits docs/.

Usage: python3 pptx2web_native.py <input.pptx> [--out docs]
"""
import argparse, html, json, os, re, shutil, sys, zipfile
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

PTW, PTH = 960.0, 540.0  # 16:9 slide in points


def pct(v, total):
    return round((v or 0) / total * 100, 3)


def color_of(font):
    try:
        if font.color and font.color.type is not None and font.color.rgb is not None:
            return "#" + str(font.color.rgb)
    except Exception:
        pass
    return None


def align_of(p):
    a = p.alignment
    return {1: "left", 2: "center", 3: "right", 4: "justify"}.get(int(a), "left") if a else "left"


def text_html(tf):
    out = []
    for p in tf.paragraphs:
        runs = []
        for r in p.runs:
            sz = r.font.size.pt if r.font.size else None
            st = f"font-size:{sz/PTH*100:.2f}cqh;" if sz else ""
            if r.font.bold:
                st += "font-weight:700;"
            if r.font.italic:
                st += "font-style:italic;"
            c = color_of(r.font)
            if c and c.lower() not in ("#000000", "#1b1b1b", "#212121"):
                st += f"color:{c};"
            runs.append(f'<span style="{st}">{html.escape(r.text)}</span>')
        if not runs:
            runs.append("<br>")
        lvl = "margin-left:%dem;" % (p.level * 1) if p.level else ""
        out.append(f'<p style="text-align:{align_of(p)};{lvl}">{"".join(runs)}</p>')
    return "".join(out)


def shape_div(sh, W, H, media):
    st = (f"left:{pct(sh.left,W)}%;top:{pct(sh.top,H)}%;"
          f"width:{pct(sh.width,W)}%;height:{pct(sh.height,H)}%;")
    if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            img = sh.image
            name = f"img_{img.sha1[:10]}.{img.ext}"
            (media/name).write_bytes(img.blob)
            return f'<div class="s" style="{st}"><img src="media/{name}"/></div>'
        except Exception:
            return ""
    if sh.has_table:
        rows = []
        for row in sh.table.rows:
            cells = "".join(f"<td>{html.escape(c.text)}</td>" for c in row.cells)
            rows.append(f"<tr>{cells}</tr>")
        return f'<div class="s" style="{st}"><table>{"".join(rows)}</table></div>'
    if sh.has_text_frame and sh.text_frame.text.strip():
        return f'<div class="s" style="{st}">{text_html(sh.text_frame)}</div>'
    return ""


def online_videos(pptx):
    out = {}
    with zipfile.ZipFile(pptx) as z:
        for nm in z.namelist():
            m = re.match(r"ppt/slides/_rels/slide(\d+)\.xml\.rels", nm)
            if not m:
                continue
            d = z.read(nm).decode("utf8", "ignore")
            urls = [u.replace("&amp;", "&") for u in re.findall(r'Target="([^"]+)"[^>]*External', d) if u.startswith("http")]
            urls = [u for u in urls if any(k in u for k in ("youtube", "youtu.be", "vimeo", "embed"))]
            if urls:
                out[int(m.group(1))-1] = urls
    return out


def slide_title(s, n):
    for sh in s.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            return re.sub(r"\s+", " ", sh.text_frame.text.strip())[:40]
    return f"슬라이드 {n}"


def slide_notes(s):
    if not s.has_notes_slide:
        return ""
    return re.sub(r"\n{3,}", "\n\n", s.notes_slide.notes_text_frame.text.strip())


def update_slides(updates_path):
    """Build extra 'latest update' slide HTML + titles from updates.json."""
    if not updates_path or not os.path.exists(updates_path):
        return [], [], []
    u = json.loads(Path(updates_path).read_text())
    src = u.get("source", ""); fetched = u.get("fetched", "")
    slides, titles = [], []
    intro = ('<div class="s" style="left:6%;top:18%;width:88%;height:64%;">'
             '<p style="font-size:6cqh;">🆕 최신 업데이트</p>'
             f'<p style="font-size:3cqh;">출처: {html.escape(src)}</p>'
             f'<p style="font-size:2.6cqh;">수집일 {html.escape(fetched)}</p></div>')
    slides.append(intro); titles.append("최신 업데이트")
    for g in u.get("groups", []):
        lis = "".join(f'<p style="font-size:2.9cqh;">• {html.escape(it)}</p>' for it in g["items"])
        slides.append(f'<div class="s" style="left:6%;top:10%;width:88%;height:82%;">'
                      f'<p style="font-size:4.4cqh;color:#22d3ee;">{html.escape(g["category"])}</p>{lis}</div>')
        titles.append("업데이트 · " + g["category"][:24])
    return slides, titles, [""] * len(slides)


def build(pptx, out, updates=None):
    out = Path(out)
    media = out/"media"
    for d in (out, media):
        d.mkdir(parents=True, exist_ok=True)
    for f in media.glob("*"):
        f.unlink()
    p = Presentation(pptx)
    W, H = p.slide_width, p.slide_height
    onv = online_videos(pptx)
    slides, titles = [], []
    for i, s in enumerate(p.slides):
        titles.append(slide_title(s, i+1))
        parts = [shape_div(sh, W, H, media) for sh in s.shapes]
        for u in onv.get(i, []):
            parts.append(f'<div class="s vid" style="left:0;top:0;width:100%;height:100%;"><iframe src="{html.escape(u)}" allow="autoplay;encrypted-media" allowfullscreen></iframe></div>')
        slides.append("".join(x for x in parts if x))
    notes = [slide_notes(s) for s in p.slides]
    us, ut, un = update_slides(updates)
    slides += us; titles += ut; notes += un
    deck = {"title": Path(pptx).stem, "aspect": W/H, "slides": slides, "titles": titles, "notes": notes}
    (out/"deck.json").write_text(json.dumps(deck, ensure_ascii=False))
    (out/"index.html").write_text(HTML.replace("__DECK__", json.dumps(deck, ensure_ascii=False)))
    return len(slides), len(onv)


HTML = r"""<!doctype html><html lang="ko"><head><meta charset="utf-8"/>
<meta name="viewport" content="width=device-width,initial-scale=1"/><title>Deck</title>
<link rel="preconnect" href="https://fonts.googleapis.com"><link rel="preconnect" href="https://fonts.gstatic.com" crossorigin>
<link href="https://fonts.googleapis.com/css2?family=Poppins:wght@600;800&family=Noto+Sans+KR:wght@400;500;700;900&display=swap" rel="stylesheet">
<style>
*{box-sizing:border-box}
html,body{margin:0;height:100%;font-family:'Noto Sans KR','Poppins',system-ui,sans-serif}
body{display:flex;background:#070a16;color:#e6edf3}
#side{width:280px;flex:0 0 280px;height:100vh;overflow-y:auto;background:#0b1020;border-right:1px solid #1c2540;padding:22px 0}
#side h1{font-family:Poppins;font-size:15px;font-weight:800;margin:0 20px 18px;line-height:1.4;background:linear-gradient(90deg,#22d3ee,#7b2ff7,#ff5ea0);-webkit-background-clip:text;background-clip:text;color:transparent}
#side a{display:flex;gap:10px;align-items:center;padding:9px 20px;color:#9fb0d0;text-decoration:none;font-size:13px;border-left:3px solid transparent;transition:.15s}
#side a:hover{background:#131a30;color:#fff}
#side a.on{background:#131a30;color:#fff;border-left-color:#22d3ee}
#side a b{font-family:Poppins;color:#445;min-width:22px}#side a.on b{color:#22d3ee}
#scroll{flex:1;height:100vh;overflow-y:auto;scroll-snap-type:y mandatory;scroll-behavior:smooth;padding:4vh 4vw}
section{scroll-snap-align:start;min-height:92vh;display:flex;align-items:center;justify-content:center;margin-bottom:4vh}
.slide{position:relative;width:100%;max-width:1280px;aspect-ratio:var(--ar,16/9);container-type:size;color:#fff;border-radius:18px;overflow:hidden;box-shadow:0 24px 80px rgba(0,0,0,.55),0 0 0 1px rgba(255,255,255,.06)}
.slide::before{content:"";position:absolute;inset:0;background:var(--bg);z-index:-2}
.slide::after{content:"";position:absolute;inset:0;z-index:-1;background:radial-gradient(60cqw 60cqw at 100% 0%,var(--a)55,transparent),radial-gradient(50cqw 50cqw at 0% 100%,var(--b)44,transparent)}
.s{position:absolute;overflow:hidden;display:flex;flex-direction:column;justify-content:center}
.s p{margin:0 0 .25em;font-size:3.4cqh;line-height:1.3;text-shadow:0 2px 12px rgba(0,0,0,.35)}
.s p:first-child:last-child{font-size:5cqh}
.s img{width:100%;height:100%;object-fit:contain;filter:drop-shadow(0 8px 24px rgba(0,0,0,.4));border-radius:10px}
.s table{width:100%;border-collapse:separate;border-spacing:0;font-size:2.6cqh;border-radius:10px;overflow:hidden}
.s td{border:1px solid rgba(255,255,255,.18);padding:.4em .6em;background:rgba(255,255,255,.06)}
.s tr:first-child td{background:rgba(255,255,255,.16);font-weight:700}
.s.vid iframe{width:100%;height:100%;border:0;border-radius:12px}
.note{width:100%;max-width:1280px;margin:18px auto 0;background:#0b1020;border:1px solid #1c2540;border-radius:14px;padding:16px 20px;color:#aebbd6;font-size:15px;line-height:1.65;white-space:pre-wrap}
.note b{display:block;font-family:Poppins;font-size:12px;letter-spacing:.08em;color:#22d3ee;margin-bottom:6px}
.note.empty{display:none}
section{flex-direction:column}
@media(max-width:760px){#side{width:64px;flex-basis:64px}#side h1,#side a span{display:none}}
</style></head><body>
<nav id="side"><h1 id="dt"></h1><div id="toc"></div></nav>
<main id="scroll"></main><script>
const D=__DECK__;
const TH=[['#0f2027','#22d3ee','#2c5364'],['#3a1c71','#ff5ea0','#d76d77'],['#0b486b','#3bb78f','#0f2027'],['#42275a','#ff5ea0','#734b6d'],['#1f1c2c','#22d3ee','#928dab'],['#16222a','#a6ed5d','#3a6073']];
dt.textContent=D.title;
scrollEl=document.getElementById('scroll');toc=document.getElementById('toc');
D.slides.forEach((h,i)=>{const t=TH[i%TH.length];
  scrollEl.insertAdjacentHTML('beforeend','<section id="s'+i+'"><div class="slide" style="--ar:'+(D.aspect||16/9)+';--bg:linear-gradient(135deg,'+t[0]+','+t[2]+');--a:'+t[1]+';--b:'+t[2]+'">'+h+'</div></section>');
  const sec=scrollEl.lastElementChild, nt=(D.notes&&D.notes[i])||'';
  const nd=document.createElement('div');nd.className='note'+(nt?'':' empty');nd.innerHTML='<b>발표 스크립트</b>';const sp=document.createElement('span');sp.textContent=nt;nd.appendChild(sp);sec.appendChild(nd);
  toc.insertAdjacentHTML('beforeend','<a href="#s'+i+'" data-i="'+i+'"><b>'+(i+1).toString().padStart(2,'0')+'</b><span>'+(D.titles[i]||'')+'</span></a>');});
const links=[...toc.querySelectorAll('a')];
const ob=new IntersectionObserver(es=>es.forEach(e=>{if(e.isIntersecting){links.forEach(l=>l.classList.remove('on'));const a=links[+e.target.id.slice(1)];a.classList.add('on');a.scrollIntoView({block:'nearest'});}}),{root:scrollEl,threshold:.5});
document.querySelectorAll('section').forEach(s=>ob.observe(s));
</script></body></html>"""


# ---------- reflow mode: screenshot + web-native reinterpretation ----------

def find_soffice():
    for c in ("soffice", "libreoffice", "/opt/homebrew/bin/soffice",
              "/Applications/LibreOffice.app/Contents/MacOS/soffice"):
        if shutil.which(c) or os.path.exists(c):
            return c
    return None


def render_pngs(pptx, outdir, scale):
    import tempfile, subprocess
    so = find_soffice()
    if not so:
        return 0
    with tempfile.TemporaryDirectory() as tmp:
        subprocess.run([so, "--headless", "--convert-to", "pdf", "--outdir", tmp, pptx],
                       check=True, capture_output=True)
        pdfs = list(Path(tmp).glob("*.pdf"))
        if not pdfs:
            return 0
        import fitz
        doc = fitz.open(pdfs[0]); n = doc.page_count
        for i in range(n):
            doc[i].get_pixmap(matrix=fitz.Matrix(scale, scale)).save(str(outdir/f"slide-{i+1:03d}.png"))
        return n


def slide_points(s):
    pts, title = [], None
    for sh in s.shapes:
        if not sh.has_text_frame:
            continue
        for p in sh.text_frame.paragraphs:
            t = re.sub(r"\s+", " ", p.text.strip())
            if not t:
                continue
            if title is None:
                title = t
            else:
                pts.append(t)
    return title, pts


def build_reflow(pptx, out, scale, updates=None):
    out = Path(out); media = out/"media"
    for d in (out, media):
        d.mkdir(parents=True, exist_ok=True)
    for f in media.glob("*"):
        f.unlink()
    n = render_pngs(pptx, media, scale)
    p = Presentation(pptx)
    onv = online_videos(pptx)
    cards = []
    for i, s in enumerate(p.slides):
        title, pts = slide_points(s)
        note = slide_notes(s)
        cards.append({
            "n": i+1,
            "title": title or f"슬라이드 {i+1}",
            "points": pts,
            "note": note,
            "img": f"media/slide-{i+1:03d}.png" if n else "",
            "videos": onv.get(i, []),
        })
    titles = [c["title"][:40] for c in cards]
    us, ut, _ = update_slides(updates)
    for j, (h, t) in enumerate(zip(us, ut)):
        cards.append({"n": len(cards)+1, "title": t, "html": h, "points": [], "note": "", "img": "", "videos": []})
        titles.append(t)
    deck = {"title": Path(pptx).stem, "cards": cards, "titles": titles}
    (out/"deck.json").write_text(json.dumps(deck, ensure_ascii=False))
    (out/"index.html").write_text(REFLOW.replace("__DECK__", json.dumps(deck, ensure_ascii=False)))
    return len(cards), sum(len(c["videos"]) for c in cards)


REFLOW = r"""<!doctype html><html lang="ko"><head><meta charset="utf-8"/>
<meta name="viewport" content="width=device-width,initial-scale=1"/><title>Deck</title>
<link rel="preconnect" href="https://fonts.googleapis.com"><link rel="preconnect" href="https://fonts.gstatic.com" crossorigin>
<link href="https://fonts.googleapis.com/css2?family=Poppins:wght@600;800&family=Noto+Sans+KR:wght@400;500;700;900&display=swap" rel="stylesheet">
<style>
*{box-sizing:border-box}
html,body{margin:0;height:100%;font-family:'Noto Sans KR','Poppins',system-ui,sans-serif}
body{display:flex;background:#070a16;color:#e6edf3}
#side{width:280px;flex:0 0 280px;height:100vh;overflow-y:auto;background:#0b1020;border-right:1px solid #1c2540;padding:22px 0}
#side h1{font-family:Poppins;font-size:15px;font-weight:800;margin:0 20px 18px;line-height:1.4;background:linear-gradient(90deg,#22d3ee,#7b2ff7,#ff5ea0);-webkit-background-clip:text;background-clip:text;color:transparent}
#side a{display:flex;gap:10px;align-items:center;padding:9px 20px;color:#9fb0d0;text-decoration:none;font-size:13px;border-left:3px solid transparent;transition:.15s}
#side a:hover{background:#131a30;color:#fff}#side a.on{background:#131a30;color:#fff;border-left-color:#22d3ee}
#side a b{font-family:Poppins;color:#445;min-width:22px}#side a.on b{color:#22d3ee}
#scroll{flex:1;height:100vh;overflow-y:auto;scroll-snap-type:y proximity;scroll-behavior:smooth;padding:5vh 5vw}
section{scroll-snap-align:start;min-height:90vh;display:flex;flex-direction:column;justify-content:center;margin-bottom:7vh;max-width:1100px;margin-inline:auto}
.t{font-family:Poppins,'Noto Sans KR';font-size:clamp(28px,4vw,52px);font-weight:800;line-height:1.2;margin:0 0 22px;background:linear-gradient(120deg,var(--a),var(--b));-webkit-background-clip:text;background-clip:text;color:transparent}
.row{display:grid;grid-template-columns:1.1fr .9fr;gap:30px;align-items:start}
.row.solo{grid-template-columns:1fr}
ul{margin:0;padding:0;list-style:none}
li{position:relative;padding:12px 0 12px 26px;font-size:clamp(16px,1.6vw,21px);line-height:1.6;border-bottom:1px solid rgba(255,255,255,.06)}
li::before{content:"";position:absolute;left:0;top:20px;width:11px;height:11px;border-radius:3px;background:linear-gradient(135deg,var(--a),var(--b))}
.shot{border-radius:14px;overflow:hidden;border:1px solid rgba(255,255,255,.1);box-shadow:0 20px 60px rgba(0,0,0,.5)}
.shot img{width:100%;display:block}
.vid iframe{width:100%;aspect-ratio:16/9;border:0;border-radius:14px}
.note{margin-top:22px;background:#0b1020;border:1px solid #1c2540;border-radius:14px;padding:16px 20px;color:#aebbd6;font-size:15px;line-height:1.7;white-space:pre-wrap}
.note b{display:block;font-family:Poppins;font-size:12px;letter-spacing:.08em;color:#22d3ee;margin-bottom:6px}
.note.empty{display:none}.solo .shot{max-width:760px;margin:0 auto}
.html p{font-size:clamp(16px,1.7vw,22px);line-height:1.5;margin:.2em 0}
@media(max-width:860px){.row{grid-template-columns:1fr}#side{width:64px;flex-basis:64px}#side h1,#side a span{display:none}}
</style></head><body>
<nav id="side"><h1 id="dt"></h1><div id="toc"></div></nav>
<main id="scroll"></main><script>
const D=__DECK__;
const TH=[['#22d3ee','#7b2ff7'],['#ff5ea0','#7b2ff7'],['#3bb78f','#0b9'],['#ff8a00','#ff5ea0'],['#22d3ee','#3bb78f'],['#a6ed5d','#22d3ee']];
dt.textContent=D.title;scrollEl=document.getElementById('scroll');toc=document.getElementById('toc');
D.cards.forEach((c,i)=>{const t=TH[i%TH.length];let body;
  if(c.html){body='<div class="html">'+c.html+'</div>';}
  else{const pts=c.points.map(x=>'<li>'+x.replace(/</g,'&lt;')+'</li>').join('');
    const vis=c.videos&&c.videos.length?'<div class="vid"><iframe src="'+c.videos[0]+'" allow="autoplay;encrypted-media" allowfullscreen></iframe></div>':(c.img?'<div class="shot"><img loading="lazy" src="'+c.img+'"></div>':'');
    const solo=pts?'':' solo';body='<div class="row'+solo+'">'+(pts?'<ul>'+pts+'</ul>':'')+vis+'</div>';}
  scrollEl.insertAdjacentHTML('beforeend','<section id="s'+i+'" style="--a:'+t[0]+';--b:'+t[1]+'"><h2 class="t">'+(c.title||'').replace(/</g,'&lt;')+'</h2>'+body+'</section>');
  const sec=scrollEl.lastElementChild;const nd=document.createElement('div');nd.className='note'+(c.note?'':' empty');nd.innerHTML='<b>발표 스크립트</b>';const sp=document.createElement('span');sp.textContent=c.note||'';nd.appendChild(sp);sec.appendChild(nd);
  toc.insertAdjacentHTML('beforeend','<a href="#s'+i+'"><b>'+(i+1).toString().padStart(2,'0')+'</b><span>'+(D.titles[i]||'')+'</span></a>');});
const links=[...toc.querySelectorAll('a')];
const ob=new IntersectionObserver(es=>es.forEach(e=>{if(e.isIntersecting){links.forEach(l=>l.classList.remove('on'));const a=links[+e.target.id.slice(1)];a.classList.add('on');a.scrollIntoView({block:'nearest'});}}),{root:scrollEl,threshold:.4});
document.querySelectorAll('section').forEach(s=>ob.observe(s));
</script></body></html>"""


if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("pptx"); ap.add_argument("--out", default="docs"); ap.add_argument("--updates", default=None)
    ap.add_argument("--reflow", action="store_true", help="web-native reinterpretation: screenshot + key points + notes")
    ap.add_argument("--scale", type=float, default=2.0)
    a = ap.parse_args()
    if not os.path.exists(a.pptx):
        sys.exit(f"Input not found: {a.pptx}")
    if a.reflow:
        n, v = build_reflow(a.pptx, a.out, a.scale, a.updates)
        print(f"OK(reflow): {n} slides, {v} online videos -> {a.out}/")
    else:
        n, v = build(a.pptx, a.out, a.updates)
        print(f"OK: {n} slides, {v} online videos -> {a.out}/")
