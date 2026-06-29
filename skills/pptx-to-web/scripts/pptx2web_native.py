#!/usr/bin/env python3
"""pptx2web_native — rebuild a .pptx as a real HTML deck (content, not screenshots).

Extracts text runs (font, size, color, bold/italic, align), pictures, tables and
embedded/online videos from each slide and lays them out as absolutely-positioned
HTML/CSS so text stays selectable and the layout adapts to the web. Emits docs/.

Usage: python3 pptx2web_native.py <input.pptx> [--out docs]
"""
import argparse, html, json, os, re, shutil, sys, zipfile
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

PTW, PTH = 960.0, 540.0  # 16:9 slide in points


def pct(v, total):
    return round((v or 0) / total * 100, 3)


def color_of(font):
    try:
        if font.color and font.color.type is not None and font.color.rgb is not None:
            return "#" + str(font.color.rgb)
    except Exception:
        pass
    return None


def align_of(p):
    a = p.alignment
    return {1: "left", 2: "center", 3: "right", 4: "justify"}.get(int(a), "left") if a else "left"


def text_html(tf):
    out = []
    for p in tf.paragraphs:
        runs = []
        for r in p.runs:
            sz = r.font.size.pt if r.font.size else None
            st = f"font-size:{sz/PTH*100:.2f}cqh;" if sz else ""
            if r.font.bold:
                st += "font-weight:700;"
            if r.font.italic:
                st += "font-style:italic;"
            c = color_of(r.font)
            if c and c.lower() not in ("#000000", "#1b1b1b", "#212121"):
                st += f"color:{c};"
            runs.append(f'<span style="{st}">{html.escape(r.text)}</span>')
        if not runs:
            runs.append("<br>")
        lvl = "margin-left:%dem;" % (p.level * 1) if p.level else ""
        out.append(f'<p style="text-align:{align_of(p)};{lvl}">{"".join(runs)}</p>')
    return "".join(out)


def shape_div(sh, W, H, media):
    st = (f"left:{pct(sh.left,W)}%;top:{pct(sh.top,H)}%;"
          f"width:{pct(sh.width,W)}%;height:{pct(sh.height,H)}%;")
    if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            img = sh.image
            name = f"img_{img.sha1[:10]}.{img.ext}"
            (media/name).write_bytes(img.blob)
            return f'<div class="s" style="{st}"><img src="media/{name}"/></div>'
        except Exception:
            return ""
    if sh.has_table:
        rows = []
        for row in sh.table.rows:
            cells = "".join(f"<td>{html.escape(c.text)}</td>" for c in row.cells)
            rows.append(f"<tr>{cells}</tr>")
        return f'<div class="s" style="{st}"><table>{"".join(rows)}</table></div>'
    if sh.has_text_frame and sh.text_frame.text.strip():
        return f'<div class="s" style="{st}">{text_html(sh.text_frame)}</div>'
    return ""


def online_videos(pptx):
    out = {}
    with zipfile.ZipFile(pptx) as z:
        for nm in z.namelist():
            m = re.match(r"ppt/slides/_rels/slide(\d+)\.xml\.rels", nm)
            if not m:
                continue
            d = z.read(nm).decode("utf8", "ignore")
            urls = [u.replace("&amp;", "&") for u in re.findall(r'Target="([^"]+)"[^>]*External', d) if u.startswith("http")]
            urls = [u for u in urls if any(k in u for k in ("youtube", "youtu.be", "vimeo", "embed"))]
            if urls:
                out[int(m.group(1))-1] = urls
    return out


def slide_title(s, n):
    for sh in s.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            return re.sub(r"\s+", " ", sh.text_frame.text.strip())[:40]
    return f"슬라이드 {n}"


def slide_notes(s):
    if not s.has_notes_slide:
        return ""
    return re.sub(r"\n{3,}", "\n\n", s.notes_slide.notes_text_frame.text.strip())


def update_slides(updates_path):
    """Build extra 'latest update' slide HTML + titles from updates.json."""
    if not updates_path or not os.path.exists(updates_path):
        return [], [], []
    u = json.loads(Path(updates_path).read_text())
    src = u.get("source", ""); fetched = u.get("fetched", "")
    slides, titles = [], []
    intro = ('<div class="s" style="left:6%;top:18%;width:88%;height:64%;">'
             '<p style="font-size:6cqh;">🆕 최신 업데이트</p>'
             f'<p style="font-size:3cqh;">출처: {html.escape(src)}</p>'
             f'<p style="font-size:2.6cqh;">수집일 {html.escape(fetched)}</p></div>')
    slides.append(intro); titles.append("최신 업데이트")
    for g in u.get("groups", []):
        lis = "".join(f'<p style="font-size:2.9cqh;">• {html.escape(it)}</p>' for it in g["items"])
        slides.append(f'<div class="s" style="left:6%;top:10%;width:88%;height:82%;">'
                      f'<p style="font-size:4.4cqh;color:#22d3ee;">{html.escape(g["category"])}</p>{lis}</div>')
        titles.append("업데이트 · " + g["category"][:24])
    return slides, titles, [""] * len(slides)


def build(pptx, out, updates=None):
    out = Path(out)
    media = out/"media"
    for d in (out, media):
        d.mkdir(parents=True, exist_ok=True)
    for f in media.glob("*"):
        f.unlink()
    p = Presentation(pptx)
    W, H = p.slide_width, p.slide_height
    onv = online_videos(pptx)
    slides, titles = [], []
    for i, s in enumerate(p.slides):
        titles.append(slide_title(s, i+1))
        parts = [shape_div(sh, W, H, media) for sh in s.shapes]
        for u in onv.get(i, []):
            parts.append(f'<div class="s vid" style="left:0;top:0;width:100%;height:100%;"><iframe src="{html.escape(u)}" allow="autoplay;encrypted-media" allowfullscreen></iframe></div>')
        slides.append("".join(x for x in parts if x))
    notes = [slide_notes(s) for s in p.slides]
    us, ut, un = update_slides(updates)
    slides += us; titles += ut; notes += un
    deck = {"title": Path(pptx).stem, "aspect": W/H, "slides": slides, "titles": titles, "notes": notes}
    (out/"deck.json").write_text(json.dumps(deck, ensure_ascii=False))
    (out/"index.html").write_text(HTML.replace("__DECK__", json.dumps(deck, ensure_ascii=False)))
    return len(slides), len(onv)


HTML = r"""<!doctype html><html lang="ko"><head><meta charset="utf-8"/>
<meta name="viewport" content="width=device-width,initial-scale=1"/><title>Deck</title>
<link rel="preconnect" href="https://fonts.googleapis.com"><link rel="preconnect" href="https://fonts.gstatic.com" crossorigin>
<link href="https://fonts.googleapis.com/css2?family=Poppins:wght@600;800&family=Noto+Sans+KR:wght@400;500;700;900&display=swap" rel="stylesheet">
<style>
*{box-sizing:border-box}
html,body{margin:0;height:100%;font-family:'Noto Sans KR','Poppins',system-ui,sans-serif}
body{display:flex;background:#070a16;color:#e6edf3}
#side{width:280px;flex:0 0 280px;height:100vh;overflow-y:auto;background:#0b1020;border-right:1px solid #1c2540;padding:22px 0}
#side h1{font-family:Poppins;font-size:15px;font-weight:800;margin:0 20px 18px;line-height:1.4;background:linear-gradient(90deg,#22d3ee,#7b2ff7,#ff5ea0);-webkit-background-clip:text;background-clip:text;color:transparent}
#side a{display:flex;gap:10px;align-items:center;padding:9px 20px;color:#9fb0d0;text-decoration:none;font-size:13px;border-left:3px solid transparent;transition:.15s}
#side a:hover{background:#131a30;color:#fff}
#side a.on{background:#131a30;color:#fff;border-left-color:#22d3ee}
#side a b{font-family:Poppins;color:#445;min-width:22px}#side a.on b{color:#22d3ee}
#scroll{flex:1;height:100vh;overflow-y:auto;scroll-snap-type:y mandatory;scroll-behavior:smooth;padding:4vh 4vw}
section{scroll-snap-align:start;min-height:92vh;display:flex;align-items:center;justify-content:center;margin-bottom:4vh}
.slide{position:relative;width:100%;max-width:1280px;aspect-ratio:var(--ar,16/9);container-type:size;color:#fff;border-radius:18px;overflow:hidden;box-shadow:0 24px 80px rgba(0,0,0,.55),0 0 0 1px rgba(255,255,255,.06)}
.slide::before{content:"";position:absolute;inset:0;background:var(--bg);z-index:-2}
.slide::after{content:"";position:absolute;inset:0;z-index:-1;background:radial-gradient(60cqw 60cqw at 100% 0%,var(--a)55,transparent),radial-gradient(50cqw 50cqw at 0% 100%,var(--b)44,transparent)}
.s{position:absolute;overflow:hidden;display:flex;flex-direction:column;justify-content:center}
.s p{margin:0 0 .25em;font-size:3.4cqh;line-height:1.3;text-shadow:0 2px 12px rgba(0,0,0,.35)}
.s p:first-child:last-child{font-size:5cqh}
.s img{width:100%;height:100%;object-fit:contain;filter:drop-shadow(0 8px 24px rgba(0,0,0,.4));border-radius:10px}
.s table{width:100%;border-collapse:separate;border-spacing:0;font-size:2.6cqh;border-radius:10px;overflow:hidden}
.s td{border:1px solid rgba(255,255,255,.18);padding:.4em .6em;background:rgba(255,255,255,.06)}
.s tr:first-child td{background:rgba(255,255,255,.16);font-weight:700}
.s.vid iframe{width:100%;height:100%;border:0;border-radius:12px}
.note{width:100%;max-width:1280px;margin:18px auto 0;background:#0b1020;border:1px solid #1c2540;border-radius:14px;padding:16px 20px;color:#aebbd6;font-size:15px;line-height:1.65;white-space:pre-wrap}
.note b{display:block;font-family:Poppins;font-size:12px;letter-spacing:.08em;color:#22d3ee;margin-bottom:6px}
.note.empty{display:none}
section{flex-direction:column}
@media(max-width:760px){#side{width:64px;flex-basis:64px}#side h1,#side a span{display:none}}
</style></head><body>
<nav id="side"><h1 id="dt"></h1><div id="toc"></div></nav>
<main id="scroll"></main><script>
const D=__DECK__;
const TH=[['#0f2027','#22d3ee','#2c5364'],['#3a1c71','#ff5ea0','#d76d77'],['#0b486b','#3bb78f','#0f2027'],['#42275a','#ff5ea0','#734b6d'],['#1f1c2c','#22d3ee','#928dab'],['#16222a','#a6ed5d','#3a6073']];
dt.textContent=D.title;
scrollEl=document.getElementById('scroll');toc=document.getElementById('toc');
D.slides.forEach((h,i)=>{const t=TH[i%TH.length];
  scrollEl.insertAdjacentHTML('beforeend','<section id="s'+i+'"><div class="slide" style="--ar:'+(D.aspect||16/9)+';--bg:linear-gradient(135deg,'+t[0]+','+t[2]+');--a:'+t[1]+';--b:'+t[2]+'">'+h+'</div></section>');
  const sec=scrollEl.lastElementChild, nt=(D.notes&&D.notes[i])||'';
  const nd=document.createElement('div');nd.className='note'+(nt?'':' empty');nd.innerHTML='<b>발표 스크립트</b>';const sp=document.createElement('span');sp.textContent=nt;nd.appendChild(sp);sec.appendChild(nd);
  toc.insertAdjacentHTML('beforeend','<a href="#s'+i+'" data-i="'+i+'"><b>'+(i+1).toString().padStart(2,'0')+'</b><span>'+(D.titles[i]||'')+'</span></a>');});
const links=[...toc.querySelectorAll('a')];
const ob=new IntersectionObserver(es=>es.forEach(e=>{if(e.isIntersecting){links.forEach(l=>l.classList.remove('on'));const a=links[+e.target.id.slice(1)];a.classList.add('on');a.scrollIntoView({block:'nearest'});}}),{root:scrollEl,threshold:.5});
document.querySelectorAll('section').forEach(s=>ob.observe(s));
</script></body></html>"""


# ---------- reflow mode: screenshot + web-native reinterpretation ----------

def find_soffice():
    for c in ("soffice", "libreoffice", "/opt/homebrew/bin/soffice",
              "/Applications/LibreOffice.app/Contents/MacOS/soffice"):
        if shutil.which(c) or os.path.exists(c):
            return c
    return None


def render_pngs(pptx, outdir, scale):
    import tempfile, subprocess
    so = find_soffice()
    if not so:
        return 0
    with tempfile.TemporaryDirectory() as tmp:
        subprocess.run([so, "--headless", "--convert-to", "pdf", "--outdir", tmp, pptx],
                       check=True, capture_output=True)
        pdfs = list(Path(tmp).glob("*.pdf"))
        if not pdfs:
            return 0
        import fitz
        doc = fitz.open(pdfs[0]); n = doc.page_count
        for i in range(n):
            doc[i].get_pixmap(matrix=fitz.Matrix(scale, scale)).save(str(outdir/f"slide-{i+1:03d}.png"))
        return n


def _walk(shapes):
    for sh in shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _walk(sh.shapes)
        else:
            yield sh


def slide_points(s):
    pts, title = [], None
    for sh in _walk(s.shapes):
        if not sh.has_text_frame:
            continue
        for p in sh.text_frame.paragraphs:
            t = re.sub(r"\s+", " ", p.text.strip())
            if not t:
                continue
            if title is None:
                title = t
            else:
                pts.append(t)
    return title, pts


def slide_pictures(s, media, W, H):
    out = []
    for sh in _walk(s.shapes):
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                wpct = (sh.width or 0) / W; hpct = (sh.height or 0) / H
                if wpct < 0.22 or hpct < 0.18:   # skip small icons/logos
                    continue
                img = sh.image
                name = f"img_{img.sha1[:10]}.{img.ext}"
                (media/name).write_bytes(img.blob)
                out.append(f"media/{name}")
            except Exception:
                pass
    return out


def build_reflow(pptx, out, scale, updates=None):
    out = Path(out); media = out/"media"
    for d in (out, media):
        d.mkdir(parents=True, exist_ok=True)
    for f in media.glob("*"):
        f.unlink()
    p = Presentation(pptx)
    W, H = p.slide_width, p.slide_height
    onv = online_videos(pptx)
    cards = []
    for i, s in enumerate(p.slides):
        title, pts = slide_points(s)
        imgs = slide_pictures(s, media, W, H)
        note = slide_notes(s)
        videos = onv.get(i, [])
        body = bool(pts or imgs or videos)
        if not body and note:   # promote notes into readable bullets
            pts = [ln.strip() for ln in re.split(r"[\n•·]|(?<=[.!?])\s", note) if len(ln.strip()) > 4]
            body = bool(pts)
        c = {"n": i+1, "title": title or "", "points": pts, "note": note, "imgs": imgs, "videos": videos}
        if not body:   # section divider: reinterpret as a clean HTML hero, never a screenshot
            c["divider"] = True
            c["title"] = title or Path(pptx).stem
        cards.append(c)
    titles = [c["title"][:40] for c in cards]
    us, ut, _ = update_slides(updates)
    for h, t in zip(us, ut):
        cards.append({"n": len(cards)+1, "title": t, "html": h, "points": [], "note": "", "imgs": [], "videos": []})
        titles.append(t)
    deck = {"title": Path(pptx).stem, "cards": cards, "titles": titles}
    (out/"deck.json").write_text(json.dumps(deck, ensure_ascii=False))
    (out/"index.html").write_text(REFLOW.replace("__DECK__", json.dumps(deck, ensure_ascii=False)))
    empties = [c["n"] for c in cards if not (c.get("html") or c.get("divider") or c["points"] or c["imgs"] or c["videos"])]
    if empties:
        print(f"WARN: empty slides: {empties}", file=sys.stderr)
    return len(cards), sum(len(c["videos"]) for c in cards)


def build_content(content_path, out, updates=None):
    """Author-driven: read a hand-written content.json (AI reinterpretation of slides)
    and render the same landscape deck shell. Avoids empty slides entirely."""
    out = Path(out); media = out/"media"
    out.mkdir(parents=True, exist_ok=True); media.mkdir(parents=True, exist_ok=True)
    data = json.loads(Path(content_path).read_text())
    title = data.get("title", "Deck")
    cards = []; titles = []
    for i, s in enumerate(data["slides"]):
        if s.get("divider"):
            cards.append({"n": i+1, "title": s["title"], "divider": True,
                          "points": [], "note": s.get("note",""), "imgs": [], "videos": []})
        else:
            cards.append({"n": i+1, "title": s.get("title",""), "html": s.get("html",""),
                          "points": [], "note": s.get("note",""), "imgs": [], "videos": []})
        titles.append((s.get("title") or "섹션")[:40])
    us, ut, _ = update_slides(updates)
    for h, t in zip(us, ut):
        cards.append({"n": len(cards)+1, "title": t, "html": h, "points": [], "note": "", "imgs": [], "videos": []})
        titles.append(t)
    deck = {"title": title, "cards": cards, "titles": titles}
    (out/"deck.json").write_text(json.dumps(deck, ensure_ascii=False))
    (out/"index.html").write_text(REFLOW.replace("__DECK__", json.dumps(deck, ensure_ascii=False)))
    empties = [c["n"] for c in cards if not (c.get("html") or c.get("divider"))]
    if empties:
        print(f"WARN: empty slides: {empties}", file=sys.stderr)
    return len(cards), 0


REFLOW = r"""<!doctype html><html lang="ko"><head><meta charset="utf-8"/>
<meta name="viewport" content="width=device-width,initial-scale=1"/><title>Deck</title>
<link rel="preconnect" href="https://fonts.googleapis.com"><link rel="preconnect" href="https://fonts.gstatic.com" crossorigin>
<link href="https://fonts.googleapis.com/css2?family=Sora:wght@600;700;800&family=Noto+Sans+KR:wght@400;500;700;900&display=swap" rel="stylesheet">
<style>
:root{--a:#22d3ee;--b:#7b2ff7}
*{box-sizing:border-box}
html,body{margin:0;height:100%;font-family:'Noto Sans KR','Sora',system-ui,sans-serif}
body{display:flex;background:#05060d;color:#e6edf3;overflow:hidden}
#bg{position:fixed;inset:0;z-index:0;pointer-events:none;background:radial-gradient(40vw 40vw at 12% 8%,rgba(34,211,238,.18),transparent),radial-gradient(45vw 45vw at 90% 18%,rgba(123,47,247,.20),transparent),radial-gradient(50vw 50vw at 70% 95%,rgba(255,94,160,.16),transparent);animation:drift 18s ease-in-out infinite alternate}
#bg::after{content:"";position:absolute;inset:0;background-image:linear-gradient(rgba(255,255,255,.025) 1px,transparent 1px),linear-gradient(90deg,rgba(255,255,255,.025) 1px,transparent 1px);background-size:54px 54px;mask:radial-gradient(80% 60% at 60% 40%,#000,transparent)}
@keyframes drift{to{transform:translateY(-4%) scale(1.05)}}
#prog{position:fixed;top:0;left:280px;right:0;height:3px;z-index:20;background:linear-gradient(90deg,#22d3ee,#7b2ff7,#ff5ea0);transform-origin:0;transform:scaleX(0)}
#side{width:280px;flex:0 0 280px;height:100vh;overflow-y:auto;background:rgba(8,11,22,.82);backdrop-filter:blur(14px);border-right:1px solid #1c2540;padding:22px 0;z-index:10}
#side h1{font-family:Sora;font-size:15px;font-weight:800;margin:0 20px 20px;line-height:1.4;background:linear-gradient(90deg,#22d3ee,#7b2ff7,#ff5ea0);-webkit-background-clip:text;background-clip:text;color:transparent}
#side a{display:flex;gap:11px;align-items:center;padding:9px 20px;color:#9fb0d0;text-decoration:none;font-size:13px;border-left:3px solid transparent;transition:.18s}
#side a:hover{background:#131a30;color:#fff}#side a.on{background:#131a30;color:#fff;border-left-color:#22d3ee}
#side a b{font-family:Sora;color:#445;min-width:22px}#side a.on b{color:#22d3ee}
#scroll{flex:1;height:100vh;overflow-y:auto;scroll-snap-type:y proximity;scroll-behavior:smooth;z-index:1}
section{scroll-snap-align:start;min-height:62vh;display:flex;flex-direction:column;justify-content:center;padding:9vh 5vw;max-width:1320px;margin-inline:auto;overflow:visible;border-bottom:1px solid rgba(255,255,255,.05)}
section:last-of-type{border-bottom:0}
.t{position:relative;font-family:Sora,'Noto Sans KR';font-size:clamp(24px,3.2vw,44px);font-weight:800;line-height:1.12;margin:0 0 26px;padding-bottom:14px;letter-spacing:-.5px;background:linear-gradient(120deg,var(--a),var(--b));-webkit-background-clip:text;background-clip:text;color:transparent}
.t::before{content:"";position:absolute;left:0;bottom:0;width:72px;height:4px;border-radius:4px;background:linear-gradient(90deg,var(--a),var(--b));box-shadow:0 0 18px var(--a)}
.kick{display:inline-flex;align-items:center;gap:8px;font-family:Sora;font-size:12px;letter-spacing:.22em;color:#7f8db0;margin:0 0 12px;text-transform:uppercase}
.kick i{width:24px;height:2px;background:linear-gradient(90deg,var(--a),var(--b));display:inline-block;border-radius:2px}
.row{display:grid;grid-template-columns:1.05fr .95fr;gap:30px;align-items:center}
.row.solo{grid-template-columns:1fr}
ul{margin:0;padding:22px 26px;list-style:none;background:rgba(255,255,255,.045);border:1px solid rgba(255,255,255,.09);border-radius:18px;backdrop-filter:blur(8px)}
ul.cols{columns:2;column-gap:30px}
.row.solo ul.cols{columns:3}
li{position:relative;padding:9px 0 9px 26px;font-size:clamp(14px,1.4vw,19px);line-height:1.5;border-bottom:1px solid rgba(255,255,255,.06);break-inside:avoid}li:last-child{border:0}
li::before{content:"";position:absolute;left:0;top:16px;width:10px;height:10px;border-radius:3px;background:linear-gradient(135deg,var(--a),var(--b));box-shadow:0 0 12px var(--a)}
.shot{border-radius:18px;overflow:hidden;border:1px solid rgba(255,255,255,.1);box-shadow:0 24px 70px rgba(0,0,0,.55);display:flex;flex-direction:column;gap:8px;padding:8px;background:rgba(255,255,255,.05);max-height:74vh}
.shot img{width:100%;max-height:70vh;object-fit:contain;display:block;border-radius:12px}
.vid iframe{width:100%;aspect-ratio:16/9;border:0;border-radius:16px;box-shadow:0 24px 70px rgba(0,0,0,.55)}
.note{margin-top:14px;background:rgba(11,16,32,.7);border:1px solid #1c2540;border-radius:14px;padding:12px 18px;color:#aebbd6;font-size:13.5px;line-height:1.55;white-space:pre-wrap;backdrop-filter:blur(8px)}
.note b{display:block;font-family:Sora;font-size:12px;letter-spacing:.1em;color:#22d3ee;margin-bottom:7px}
.note.empty{display:none}.solo .shot{max-width:780px;margin:0 auto}
.html{background:linear-gradient(180deg,rgba(255,255,255,.06),rgba(255,255,255,.03));border:1px solid rgba(255,255,255,.1);border-radius:22px;padding:30px 34px;box-shadow:0 30px 80px rgba(0,0,0,.45),inset 0 1px 0 rgba(255,255,255,.06)}
.html p{font-size:clamp(16px,1.7vw,22px);line-height:1.55;margin:.25em 0}
.gr{display:grid;gap:18px}.gr2{grid-template-columns:1fr 1fr}.gr3{grid-template-columns:repeat(3,1fr)}.gr4{grid-template-columns:repeat(4,1fr)}
.cd{position:relative;background:linear-gradient(180deg,rgba(255,255,255,.07),rgba(255,255,255,.03));border:1px solid rgba(255,255,255,.12);border-radius:18px;padding:22px 22px;backdrop-filter:blur(8px);overflow:hidden;transition:transform .25s,box-shadow .25s,border-color .25s}
.cd::before{content:"";position:absolute;left:0;top:0;height:3px;width:100%;background:linear-gradient(90deg,var(--a),var(--b));opacity:.85}
.cd:hover{transform:translateY(-4px);border-color:rgba(255,255,255,.22);box-shadow:0 22px 50px rgba(0,0,0,.45)}
.cd h3{margin:0 0 8px;font-family:Sora,'Noto Sans KR';font-size:clamp(15px,1.5vw,20px);font-weight:800}.cd p,.cd li{font-size:clamp(13px,1.25vw,16px);line-height:1.5;color:#cdd8ee;margin:.2em 0}
.ic{width:50px;height:50px;border-radius:14px;display:flex;align-items:center;justify-content:center;font-size:24px;margin-bottom:12px;background:linear-gradient(135deg,var(--a),var(--b));box-shadow:0 10px 26px rgba(34,211,238,.32);border:1px solid rgba(255,255,255,.18)}
.flow{display:flex;align-items:stretch;gap:22px;flex-wrap:wrap}.flow .cd{flex:1;min-width:180px;position:relative}.flow .cd:not(:last-child)::after{content:"→";position:absolute;right:-17px;top:50%;transform:translateY(-50%);color:var(--a);font-size:24px;z-index:2;text-shadow:0 0 12px var(--a)}
.stat{font-family:Sora;font-size:clamp(30px,5vw,64px);font-weight:800;background:linear-gradient(120deg,var(--a),var(--b));-webkit-background-clip:text;background-clip:text;color:transparent;line-height:1}
.pill{display:inline-block;padding:5px 13px;border-radius:999px;font-size:13px;font-weight:700;margin:3px;background:linear-gradient(135deg,var(--a),var(--b));color:#03121a}
.tag{display:inline-block;padding:3px 9px;border-radius:7px;font-size:11px;font-weight:700;background:rgba(34,211,238,.16);color:#22d3ee;margin-left:6px}
.html .lead{font-size:clamp(16px,1.7vw,21px);color:#aebbd6;margin-bottom:18px}
section.divider{min-height:80vh;align-items:center}
.dvr{text-align:center}.dvr span{font-family:Sora;letter-spacing:.35em;font-size:14px;color:#7f8db0}
.dvr h2{font-family:Sora,'Noto Sans KR';font-size:clamp(40px,7vw,90px);font-weight:800;margin:.2em 0 0;background:linear-gradient(120deg,var(--a),var(--b));-webkit-background-clip:text;background-clip:text;color:transparent}
.t,.html,.row,.dvr{opacity:0;transform:translateY(18px);transition:opacity .6s ease,transform .6s ease}
section.in .t,section.in .html,section.in .row,section.in .dvr{opacity:1;transform:none}
section.in .t{transition-delay:.05s}section.in .html,section.in .row{transition-delay:.14s}
@media(max-width:860px){.row{grid-template-columns:1fr}#side{width:62px;flex-basis:62px}#side h1,#side a span{display:none}#prog{left:62px}}
</style></head><body>
<div id="bg"></div><div id="prog"></div>
<nav id="side"><h1 id="dt"></h1><div id="toc"></div></nav>
<main id="scroll"></main>
<script>
const D=__DECK__;
const TH=[['#22d3ee','#7b2ff7'],['#ff5ea0','#7b2ff7'],['#3bb78f','#22d3ee'],['#ff8a00','#ff5ea0'],['#22d3ee','#3bb78f'],['#a6ed5d','#22d3ee']];
dt.textContent=D.title;scrollEl=document.getElementById('scroll');toc=document.getElementById('toc');
D.cards.forEach((c,i)=>{const t=TH[i%TH.length];let body;
  if(c.divider){scrollEl.insertAdjacentHTML('beforeend','<section id="s'+i+'" class="divider" style="--a:'+t[0]+';--b:'+t[1]+'"><div class="dvr"><span>SECTION</span><h2>'+(c.title||'').replace(/</g,'&lt;')+'</h2></div></section>');
    toc.insertAdjacentHTML('beforeend','<a href="#s'+i+'"><b>'+(i+1).toString().padStart(2,'0')+'</b><span>'+((D.titles[i]||'섹션')).replace(/</g,'&lt;')+'</span></a>');return;}
  if(c.html){body='<div class="html">'+c.html+'</div>';}
  else{const pts=c.points.map((x,k)=>'<li>'+x.replace(/</g,'&lt;')+'</li>').join('');
    const cols=c.points.length>7?' cols':'';
    const vis=c.videos&&c.videos.length?'<div class="vid"><iframe src="'+c.videos[0]+'" allow="autoplay;encrypted-media" allowfullscreen></iframe></div>':(c.imgs&&c.imgs.length?'<div class="shot">'+c.imgs.map(s=>'<img loading="lazy" src="'+s+'">').join('')+'</div>':'');
    const solo=pts?'':' solo';body='<div class="row'+solo+'">'+(pts?'<ul class="lst'+cols+'">'+pts+'</ul>':'')+vis+'</div>';}
  scrollEl.insertAdjacentHTML('beforeend','<section id="s'+i+'" style="--a:'+t[0]+';--b:'+t[1]+'"><p class="kick"><i></i>'+(i+1).toString().padStart(2,'0')+' · '+D.title.replace(/</g,'&lt;')+'</p><h2 class="t">'+(c.title||'').replace(/</g,'&lt;')+'</h2>'+body+'</section>');
  const sec=scrollEl.lastElementChild;const nd=document.createElement('div');nd.className='note'+(c.note?'':' empty');nd.innerHTML='<b>발표 스크립트</b>';const sp=document.createElement('span');sp.textContent=c.note||'';nd.appendChild(sp);sec.appendChild(nd);
  toc.insertAdjacentHTML('beforeend','<a href="#s'+i+'"><b>'+(i+1).toString().padStart(2,'0')+'</b><span>'+((D.titles[i]||'섹션')).replace(/</g,'&lt;')+'</span></a>');});
const links=[...toc.querySelectorAll('a')],prog=document.getElementById('prog');
const ob=new IntersectionObserver(es=>es.forEach(e=>{e.target.classList.toggle('in',e.isIntersecting);if(e.isIntersecting){links.forEach(l=>l.classList.remove('on'));const a=links[+e.target.id.slice(1)];a.classList.add('on');a.scrollIntoView({block:'nearest'});}}),{root:scrollEl,threshold:.4});
document.querySelectorAll('section').forEach(s=>ob.observe(s));
scrollEl.addEventListener('scroll',()=>{prog.style.transform='scaleX('+(scrollEl.scrollTop/(scrollEl.scrollHeight-scrollEl.clientHeight||1))+')';});
</script></body></html>"""


if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("pptx", nargs="?"); ap.add_argument("--out", default="docs"); ap.add_argument("--updates", default=None)
    ap.add_argument("--reflow", action="store_true", help="web-native reinterpretation: screenshot + key points + notes")
    ap.add_argument("--content", default=None, help="author-driven JSON of AI-reinterpreted slides")
    ap.add_argument("--scale", type=float, default=2.0)
    a = ap.parse_args()
    if a.content:
        n, v = build_content(a.content, a.out, a.updates)
        print(f"OK(content): {n} slides -> {a.out}/")
    elif not os.path.exists(a.pptx):
        sys.exit(f"Input not found: {a.pptx}")
    elif a.reflow:
        n, v = build_reflow(a.pptx, a.out, a.scale, a.updates)
        print(f"OK(reflow): {n} slides, {v} online videos -> {a.out}/")
    else:
        n, v = build(a.pptx, a.out, a.updates)
        print(f"OK: {n} slides, {v} online videos -> {a.out}/")
